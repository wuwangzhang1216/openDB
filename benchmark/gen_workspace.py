#!/usr/bin/env python3
"""
Generate benchmark workspace files and upload to FileDB.

Usage:
    pip install httpx
    python gen_workspace.py [--url http://localhost:8000] [--skip-upload]
"""

import asyncio
import argparse
import json
import os
import sys

import httpx

WORKSPACE_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "benchmark_workspace")
FILEDB_URL = "http://localhost:8000"

# Map base filenames to output format (pdf/docx/pptx/csv)
FORMAT_MAP = {
    "engineering-roadmap-2025": "docx",
    "engineering-incident-report-q4": "pdf",
    "engineering-team-metrics": "csv",
    "engineering-budget-proposal-2025": "pdf",
    "engineering-architecture-review": "docx",
    "sales-q4-report": "pdf",
    "sales-pipeline-analysis": "docx",
    "sales-commission-structure": "csv",
    "sales-strategy-2025": "pptx",
    "hr-employee-census-2025": "docx",
    "hr-hiring-plan-2025": "docx",
    "hr-benefits-review": "pdf",
    "hr-employee-satisfaction-survey": "csv",
    "finance-budget-summary-2025": "pdf",
    "finance-revenue-forecast-2025": "pdf",
    "finance-expense-report-q4": "csv",
    "finance-cash-flow-analysis": "docx",
    "legal-compliance-review": "pdf",
    "legal-contract-summary-q4": "pdf",
    "legal-ip-portfolio-review": "docx",
    "executive-board-presentation-q4": "pptx",
    "risk-register-2025": "pdf",
    "strategic-priorities-2025": "pptx",
    "all-hands-meeting-notes-jan2025": "docx",
    "cross-department-project-tracker": "csv",
}

# ============================================================
# File contents — 25 files across 6 departments
# ============================================================

FILES = {}

# --- Engineering (5 files) ---

FILES["engineering-roadmap-2025.txt"] = {
    "tags": ["engineering", "2025", "roadmap"],
    "dept": "engineering",
    "content": """\
Engineering Roadmap 2025
========================

Prepared by: VP Engineering, Sarah Chen
Last Updated: January 2025

Executive Summary
-----------------
This document outlines the engineering organization's key initiatives and
milestones for fiscal year 2025. Our focus areas are platform modernization,
API expansion, mobile experience, and AI integration.

Q1 2025: Platform Migration (January - March)
----------------------------------------------
- Complete migration from monolithic architecture to microservices
- Target: 95% of traffic served by new platform by end of Q1
- Infrastructure budget: $3.2M for cloud resources
- Key dependency: hiring 15 additional backend engineers (see hr-hiring-plan-2025.txt)
- Risk: timeline pressure if hiring targets are not met by February

Q2 2025: API v3 Launch (April - June)
--------------------------------------
- Launch API v3 with GraphQL support alongside REST
- Partner integration program: onboard 20 enterprise partners
- Performance target: p99 latency under 200ms
- Budget allocation: $1.8M for developer tools and documentation
- Q4 revenue target dependency: API v3 is critical for enterprise deals in H2

Q3 2025: Mobile Redesign (July - September)
--------------------------------------------
- Complete redesign of iOS and Android applications
- Target: 4.5+ star rating on both app stores
- Budget: $2.5M for mobile engineering team expansion
- Cross-reference: see sales-strategy-2025.txt for mobile revenue projections

Q4 2025: AI Integration (October - December)
----------------------------------------------
- Integrate ML-powered recommendations into core product
- Launch AI assistant feature for enterprise tier
- Budget: $4.5M for ML infrastructure and talent
- Risk: competitive pressure from rivals who launched AI features in 2024

Total Engineering Budget Request: $12M
See finance-budget-summary-2025.txt for approval status.

Headcount Plan
--------------
Current engineering headcount: 187
Planned new hires: 45 engineers across all levels
See hr-hiring-plan-2025.txt for detailed hiring timeline.

Key Risks
---------
1. Platform migration delays could cascade to Q2 API launch
2. Hiring competition in ML/AI talent market
3. Technical debt in legacy payment system (see engineering-architecture-review.txt)
4. Budget cuts could reduce Q4 AI investment
""",
}

FILES["engineering-incident-report-q4.txt"] = {
    "tags": ["engineering", "2024", "incident"],
    "dept": "engineering",
    "content": """\
Engineering Incident Report - Q4 2024
======================================

Incident ID: INC-2024-0847
Severity: P1 (Critical)
Duration: 4 hours 23 minutes
Date: November 15, 2024

Summary
-------
A cascading failure in the payment processing service caused a complete
outage of checkout functionality from 14:07 to 18:30 UTC. The incident
affected approximately 12,000 customers and resulted in an estimated
revenue impact of $2.3M in lost transactions.

Root Cause
----------
A database connection pool exhaustion in the payment-gateway service was
triggered by a surge in retry attempts following a brief network partition
between availability zones. The connection pool was configured with a
maximum of 50 connections, which proved insufficient under the retry storm.

Timeline
--------
14:07 - First alerts triggered for elevated error rates
14:15 - On-call engineer acknowledged, began investigation
14:32 - Identified connection pool exhaustion
14:45 - Attempted pool size increase, required service restart
15:00 - Service restart failed due to pending transactions
15:30 - Escalated to VP Engineering
16:00 - Decision to implement emergency circuit breaker
17:15 - Circuit breaker deployed, partial recovery
18:30 - Full recovery confirmed

Revenue Impact
--------------
- Direct lost transactions: $2.3M
- Customer compensation (credits/refunds): $180K
- Estimated customer churn impact: $500K annually
- Total Q4 revenue impact: approximately $2.98M
- This incident reduced Q4 revenue from projected $49.3M to actual $47M
  (see sales-q4-report.txt for final Q4 revenue figures)

Action Items
------------
1. Increase connection pool limits across all services
2. Implement circuit breaker pattern in all external service calls
3. Add connection pool monitoring and alerting
4. Conduct chaos engineering exercises quarterly
5. Update risk-register-2025.txt with infrastructure resilience risk

Lessons Learned
---------------
The incident exposed a systemic risk in our infrastructure resilience.
Our disaster recovery procedures need significant improvement. This has
been added as a top priority in the engineering-architecture-review.txt.
""",
}

FILES["engineering-team-metrics.csv"] = {
    "tags": ["engineering", "2024", "metrics"],
    "dept": "engineering",
    "content": """\
Team,Headcount,Velocity,BugsFixed,Uptime,OnCallIncidents,DeployFrequency,CodeCoverage
Platform,32,87,145,99.95%,12,daily,82%
API,28,92,98,99.99%,3,daily,91%
Mobile-iOS,18,65,67,99.90%,8,weekly,74%
Mobile-Android,16,71,82,99.88%,11,weekly,71%
Data-Pipeline,22,78,53,99.97%,5,daily,88%
ML-Infrastructure,15,45,31,99.92%,7,weekly,67%
DevOps-SRE,24,N/A,42,99.96%,28,continuous,85%
Security,12,N/A,89,N/A,15,weekly,79%
Frontend,20,83,112,N/A,2,daily,76%
""",
}

FILES["engineering-budget-proposal-2025.txt"] = {
    "tags": ["engineering", "2025", "budget"],
    "dept": "engineering",
    "content": """\
Engineering Budget Proposal - Fiscal Year 2025
================================================

Submitted by: Sarah Chen, VP Engineering
Date: December 2024

Overview
--------
This proposal requests a total engineering budget of $12M for FY2025,
representing a 15% increase over FY2024 spending of $10.4M.

Budget Breakdown
----------------
1. Cloud Infrastructure: $4.2M
   - AWS/GCP compute and storage: $3.0M
   - CDN and edge services: $0.5M
   - Monitoring and observability: $0.4M
   - Security tools: $0.3M

2. Personnel (new hires only): $4.5M
   - 45 new engineering positions (see hr-hiring-plan-2025.txt)
   - Average fully-loaded cost: $100K per engineer for partial year

3. Software Licenses and Tools: $1.8M
   - Developer tooling: $0.8M
   - Third-party APIs: $0.6M
   - Testing infrastructure: $0.4M

4. ML/AI Investment: $1.5M
   - GPU compute for training: $0.8M
   - ML platform licenses: $0.4M
   - External data sources: $0.3M

Risk: Budget Cuts
-----------------
If the company-wide budget cuts discussed in finance-budget-summary-2025.txt
affect engineering, we propose the following priority tiers:

Tier 1 (must-have): Platform migration + API v3 = $7M
Tier 2 (should-have): Mobile redesign = $2.5M
Tier 3 (nice-to-have): AI integration = $2.5M

A budget cut to engineering would jeopardize the Q4 revenue target of $58M
for 2025, as the AI features are expected to drive $15M in new enterprise
revenue (see finance-revenue-forecast-2025.txt).

Comparison to Industry
----------------------
Our engineering spend as a percentage of revenue (5.7%) is below the
industry median of 8.2% for SaaS companies of our size. Increasing
investment in engineering is critical for maintaining competitive position.
""",
}

FILES["engineering-architecture-review.txt"] = {
    "tags": ["engineering", "2025", "architecture"],
    "dept": "engineering",
    "content": """\
Engineering Architecture Review - January 2025
================================================

Prepared by: Architecture Review Board

Current State Assessment
------------------------
Our technology stack consists of a primary monolithic application (Python/Django)
with several satellite microservices. The monolith handles approximately 60%
of all traffic, while microservices handle the remaining 40%.

Technical Debt Inventory
------------------------
1. Payment System (Critical)
   - Legacy payment processing module uses deprecated APIs
   - No circuit breaker implementation (see engineering-incident-report-q4.txt)
   - Estimated remediation effort: 3 engineering-months

2. Data Pipeline (High)
   - ETL jobs running on outdated Spark 2.x
   - Migration to Spark 3.x required for performance improvements
   - Estimated effort: 2 engineering-months

3. Authentication System (Medium)
   - Custom OAuth implementation should be replaced with standard library
   - Security audit flagged 3 medium-severity issues
   - Estimated effort: 1.5 engineering-months

4. Frontend Build System (Low)
   - Webpack 4 should be upgraded to modern bundler
   - Estimated effort: 1 engineering-month

Microservices Migration Plan
----------------------------
The Q1 2025 platform migration (see engineering-roadmap-2025.txt) will
decompose the monolith into the following services:

- user-service: Authentication and user management
- order-service: Order processing and fulfillment
- payment-service: Payment processing with circuit breakers
- catalog-service: Product catalog and search
- notification-service: Email, SMS, push notifications

Data Handling and Compliance
----------------------------
All services must comply with GDPR data handling requirements.
See legal-compliance-review.txt for specific requirements.

Personal data must be encrypted at rest and in transit. Data retention
policies must be implemented in all services that store user data.

Risk Assessment
---------------
The migration carries significant risk of service disruption. We recommend
a parallel-run approach where the monolith and microservices operate
simultaneously for 4-6 weeks before cutover.

See risk-register-2025.txt for the complete risk assessment.
""",
}

# --- Sales (4 files) ---

FILES["sales-q4-report.txt"] = {
    "tags": ["sales", "2024", "quarterly"],
    "dept": "sales",
    "content": """\
Sales Q4 2024 Report
=====================

Prepared by: James Morrison, VP Sales
Date: January 10, 2025

Executive Summary
-----------------
Q4 2024 was a record quarter for the sales organization. We achieved
$47M in total revenue, exceeding our Q4 revenue target of $42M by 12%.
This represents a 28% year-over-year increase from Q4 2023's $36.7M.

The outperformance was driven primarily by strong enterprise deal closings
in North America and accelerated growth in the Asia Pacific region.

Revenue Breakdown by Region
----------------------------
- North America: $26.3M (56% of total, +25% YoY)
- Asia Pacific: $13.2M (28% of total, +38% YoY)
- Europe: $7.5M (16% of total, +18% YoY)

Revenue Breakdown by Segment
------------------------------
- Enterprise: $30.1M (64%, +31% YoY)
- Mid-Market: $10.8M (23%, +22% YoY)
- SMB: $6.1M (13%, +19% YoY)

Key Deal Highlights
-------------------
1. GlobalCorp Enterprise Agreement: $4.2M (3-year)
2. TechVentures Platform License: $2.8M (2-year)
3. AsiaBank Digital Transformation: $3.5M (2-year)
4. EuroManufacturing IoT Platform: $1.9M (1-year)

The engineering incident in November (see engineering-incident-report-q4.txt)
resulted in approximately $2.3M in lost transactions, which was partially
recovered through the deals closed in December.

Q4 Revenue Target Context
--------------------------
The original Q4 revenue target of $42M was set in September 2024 based on
pipeline analysis (see sales-pipeline-analysis.txt). The target was considered
aggressive at the time given the competitive landscape. Achieving $47M
validates our enterprise-focused strategy.

FY2024 Full Year Summary
--------------------------
- Total FY2024 Revenue: $168M (+23% YoY)
- Original target: $155M (exceeded by 8.4%)
- Enterprise segment grew from 58% to 64% of revenue mix

Looking Ahead to 2025
----------------------
Based on current pipeline and market conditions, the FY2025 revenue
forecast is $210M (see finance-revenue-forecast-2025.txt). The Q4 2025
revenue target has been set at $58M, reflecting continued momentum.

See sales-strategy-2025.txt for the detailed 2025 go-to-market plan.
See legal-contract-summary-q4.txt for contract details on major deals.
""",
}

FILES["sales-pipeline-analysis.txt"] = {
    "tags": ["sales", "2024", "pipeline"],
    "dept": "sales",
    "content": """\
Sales Pipeline Analysis - Q4 2024
===================================

Prepared by: Sales Operations Team
Date: September 2024

Pipeline Overview
-----------------
This analysis informed the Q4 revenue target of $42M. Total weighted
pipeline entering Q4 was $78M across 340 opportunities.

Pipeline by Stage
-----------------
Stage 1 (Prospect): $22M across 120 opportunities (10% probability)
Stage 2 (Qualified): $18M across 95 opportunities (25% probability)
Stage 3 (Proposal): $15M across 62 opportunities (50% probability)
Stage 4 (Negotiation): $12M across 38 opportunities (75% probability)
Stage 5 (Closing): $11M across 25 opportunities (90% probability)

Weighted pipeline value: $42.3M → rounded to $42M Q4 revenue target

Pipeline by Region
------------------
North America: $38M (49% of pipeline)
Asia Pacific: $24M (31% of pipeline)
Europe: $16M (20% of pipeline)

Key Risks to Pipeline
---------------------
1. Three large enterprise deals ($8M combined) dependent on API v3
   availability (see engineering-roadmap-2025.txt for API v3 timeline)
2. APAC deals subject to currency fluctuation risk
3. European regulatory requirements may delay two deals
   (see legal-compliance-review.txt)
4. Competitive displacement risk in mid-market segment

Win Rate Analysis
-----------------
- Overall win rate: 32% (industry average: 27%)
- Enterprise win rate: 28% (improving from 24% in Q3)
- SMB win rate: 45% (stable)

Average deal size increased 18% YoY to $138K, driven by enterprise
upsell motions and platform bundle pricing.

Pipeline Velocity
-----------------
Average days to close: 72 days (down from 85 days in Q3)
Enterprise: 95 days average
Mid-Market: 55 days average
SMB: 28 days average
""",
}

FILES["sales-commission-structure.csv"] = {
    "tags": ["sales", "2024", "commission"],
    "dept": "sales",
    "content": """\
Rep,Region,Segment,Q4Quota,Q4Actual,Attainment,CommissionRate,CommissionEarned,Accelerator
Alice Wang,North America,Enterprise,3500000,4200000,120%,12%,504000,Yes
Bob Martinez,North America,Enterprise,3200000,3100000,97%,10%,310000,No
Carol Smith,North America,Mid-Market,2000000,2300000,115%,11%,253000,Yes
David Kim,Asia Pacific,Enterprise,2800000,3500000,125%,12%,420000,Yes
Emily Chen,Asia Pacific,Enterprise,2500000,2600000,104%,10%,260000,Yes
Frank Johnson,Europe,Enterprise,2200000,1900000,86%,8%,152000,No
Grace Lee,Asia Pacific,Mid-Market,1800000,2100000,117%,11%,231000,Yes
Henry Brown,North America,SMB,1200000,1400000,117%,9%,126000,Yes
Iris Patel,Europe,Mid-Market,1500000,1600000,107%,10%,160000,Yes
Jack Wilson,North America,Enterprise,3000000,2800000,93%,10%,280000,No
Kate Taylor,Europe,Enterprise,2400000,2500000,104%,10%,250000,Yes
Liam O'Brien,North America,SMB,1100000,1300000,118%,9%,117000,Yes
Maria Garcia,Asia Pacific,SMB,900000,1000000,111%,9%,90000,Yes
Nathan Park,North America,Mid-Market,1800000,2000000,111%,10%,200000,Yes
Olivia Davis,Europe,SMB,800000,850000,106%,9%,76500,Yes
""",
}

FILES["sales-strategy-2025.txt"] = {
    "tags": ["sales", "2025", "strategy"],
    "dept": "sales",
    "content": """\
Sales Strategy 2025
====================

Prepared by: James Morrison, VP Sales
Date: January 2025

Strategic Objectives
--------------------
1. Achieve $210M in annual revenue (25% YoY growth)
2. Expand enterprise segment to 70% of revenue mix
3. Enter 3 new geographic markets (Southeast Asia, Middle East, Latin America)
4. Launch partner channel contributing 15% of new bookings

Revenue Targets by Quarter
---------------------------
Q1 2025: $48M (platform migration enables enterprise features)
Q2 2025: $50M (API v3 launch drives partner revenue)
Q3 2025: $54M (mobile redesign captures consumer market)
Q4 2025: $58M (AI features drive enterprise upsell)

Total: $210M

Budget Requirements
-------------------
Sales organization budget request: $28M
- Headcount: 35 new sales hires (see hr-hiring-plan-2025.txt)
- Marketing programs: $8M
- Sales tools and enablement: $3M
- Travel and events: $2M
- Partner program: $1.5M

Note: budget cuts discussed in finance-budget-summary-2025.txt may affect
the marketing programs budget. We request that sales hiring not be reduced
as it directly impacts revenue capacity.

Go-to-Market Changes
---------------------
1. Enterprise-first: Dedicated enterprise team with specialized SE support
2. Product-led growth: Free tier funnel for SMB segment
3. Partner ecosystem: System integrator and reseller partnerships
4. Geographic expansion: Local sales teams in Singapore, Dubai, Sao Paulo

Competitive Landscape
---------------------
- Competitor A launched AI features in Q3 2024 (risk to our enterprise deals)
- Competitor B acquired a mobile-first startup (threat to mobile strategy)
- Market consolidation expected; 2 smaller competitors likely acquisition targets

Dependencies
------------
- Engineering API v3 delivery in Q2 (see engineering-roadmap-2025.txt)
- AI features in Q4 drive $15M of the $58M Q4 revenue target
- Hiring plan execution (see hr-hiring-plan-2025.txt)
""",
}

# --- HR (4 files) ---

FILES["hr-employee-census-2025.txt"] = {
    "tags": ["hr", "2025", "headcount"],
    "dept": "hr",
    "content": """\
HR Employee Census - January 2025
===================================

Prepared by: Maria Santos, CHRO
Date: January 5, 2025

Company-Wide Headcount Summary
-------------------------------
Total Active Employees: 847
- Full-time: 812
- Part-time: 23
- Contractors (long-term): 12

Headcount by Department
------------------------
Engineering:        187 (22.1%)
Sales:              156 (18.4%)
Marketing:           78 (9.2%)
Customer Success:    95 (11.2%)
Product:             54 (6.4%)
Finance:             42 (5.0%)
HR:                  28 (3.3%)
Legal:               18 (2.1%)
Operations:          65 (7.7%)
Executive:           12 (1.4%)
IT/Infrastructure:   45 (5.3%)
Data/Analytics:      38 (4.5%)
Other:               29 (3.4%)

Headcount by Region
--------------------
North America (HQ): 523 (61.7%)
Asia Pacific:       168 (19.8%)
Europe:             112 (13.2%)
Other:               44 (5.2%)

Attrition Analysis
------------------
- Annual attrition rate: 12% (industry average: 15%)
- Engineering attrition: 14% (highest department, mainly ML/AI talent)
- Sales attrition: 11%
- Voluntary attrition: 9%
- Involuntary attrition: 3%

Key concern: Engineering attrition in ML/AI roles is 22%, significantly
above the department average. This is a risk to the AI integration planned
for Q4 2025 (see engineering-roadmap-2025.txt).

Diversity Metrics
-----------------
- Gender: 42% female, 56% male, 2% non-binary
- Engineering gender split: 31% female (up from 28% in 2024)
- Management gender split: 38% female

Tenure Distribution
-------------------
- Less than 1 year: 23%
- 1-3 years: 41%
- 3-5 years: 22%
- More than 5 years: 14%

See hr-hiring-plan-2025.txt for planned headcount growth.
See hr-employee-satisfaction-survey.csv for satisfaction data.
""",
}

FILES["hr-hiring-plan-2025.txt"] = {
    "tags": ["hr", "2025", "hiring"],
    "dept": "hr",
    "content": """\
HR Hiring Plan - Fiscal Year 2025
===================================

Prepared by: Maria Santos, CHRO
Date: January 2025

Overview
--------
Plan to hire 120 new employees in FY2025, growing total headcount from
847 to approximately 967 (accounting for projected 12% attrition of ~102
departures, net growth of approximately 18 employees above replacement).

Hiring by Department
---------------------
Engineering:     45 positions
  - Backend engineers: 15
  - ML/AI engineers: 10
  - Mobile developers: 8
  - DevOps/SRE: 5
  - Frontend: 5
  - Security: 2

Sales:           35 positions
  - Enterprise AEs: 12
  - Mid-Market AEs: 8
  - SDRs/BDRs: 10
  - Sales Engineers: 5

Marketing:       20 positions
  - Content marketing: 5
  - Product marketing: 5
  - Demand generation: 5
  - Brand/creative: 5

Other:           20 positions
  - Customer Success: 8
  - Product Management: 5
  - Finance/Accounting: 3
  - Legal: 2
  - HR: 2

Hiring Timeline
---------------
Q1: 35 hires (front-loaded for platform migration support)
Q2: 30 hires
Q3: 30 hires
Q4: 25 hires

Budget Impact
-------------
- Total hiring cost (recruitment): $2.4M
- Total first-year compensation (prorated): $9.6M
- Average fully-loaded cost per hire: $120K (annualized)
- See finance-budget-summary-2025.txt for overall budget context

Risk: Budget Cuts Impact on Hiring
-----------------------------------
The finance team has proposed $8M in budget cuts across all departments
(see finance-budget-summary-2025.txt). If hiring is affected:

Scenario A (10% hiring cut): Reduce to 108 hires, defer Q3/Q4 positions
Scenario B (20% hiring cut): Reduce to 96 hires, eliminate marketing hires
Scenario C (30% hiring cut): Reduce to 84 hires, critical impact on
  engineering and sales capacity

Recommendation: Protect engineering and sales hiring as these directly
impact revenue targets. See sales-strategy-2025.txt and
engineering-roadmap-2025.txt for revenue dependencies on headcount.
""",
}

FILES["hr-benefits-review.txt"] = {
    "tags": ["hr", "2025", "benefits"],
    "dept": "hr",
    "content": """\
HR Benefits Review - FY2025
=============================

Prepared by: Benefits Team
Date: December 2024

Current Benefits Cost
---------------------
Total annual benefits cost: $18.2M (for 847 employees)
Per-employee average: $21,500/year

Benefits Breakdown
------------------
1. Health Insurance (Medical/Dental/Vision): $9.8M
   - Medical: $7.2M
   - Dental: $1.4M
   - Vision: $1.2M

2. 401(k) Match: $3.6M (6% match)

3. PTO and Leave: $2.1M (estimated cost of unused PTO liability)

4. Other Benefits: $2.7M
   - Life insurance: $0.8M
   - Disability insurance: $0.6M
   - Employee wellness: $0.5M
   - Education reimbursement: $0.4M
   - Commuter benefits: $0.4M

Proposed Budget Cuts to Benefits
---------------------------------
In response to the company-wide budget cuts (see finance-budget-summary-2025.txt),
the following changes are proposed:

1. Dental plan: Switch from Premium to Standard tier
   Savings: $420K/year
   Impact: Higher out-of-pocket costs for employees

2. Education reimbursement: Reduce from $5K to $3K per employee
   Savings: $180K/year
   Impact: May affect retention of junior engineers

3. Wellness program: Consolidate vendors
   Savings: $150K/year
   Impact: Minimal

Total proposed benefits savings: $750K/year

Note: Any changes to health benefits must comply with regulatory requirements.
See legal-compliance-review.txt for ERISA and ACA compliance obligations.

Employee Impact Assessment
--------------------------
Based on hr-employee-satisfaction-survey.csv data, benefits satisfaction
scores are already below industry benchmark (3.2/5 vs 3.8/5). Further
cuts risk increasing attrition above the current 12% rate.

Risk: If attrition increases by even 2 percentage points, the additional
replacement hiring cost would exceed the benefits savings.
""",
}

FILES["hr-employee-satisfaction-survey.csv"] = {
    "tags": ["hr", "2024", "survey"],
    "dept": "hr",
    "content": """\
Department,OverallSatisfaction,CompSatisfaction,BenefitsSatisfaction,WorkLifeBalance,CareerGrowth,ManagerRating,ResponseRate,TopConcern,RetentionRisk
Engineering,3.4,3.1,3.0,2.8,3.5,3.7,82%,Work-life balance,High
Sales,3.6,3.8,3.3,3.0,3.4,3.5,71%,Commission structure,Medium
Marketing,3.8,3.2,3.4,3.5,3.1,3.9,88%,Career growth,Medium
Customer Success,3.2,2.9,3.1,2.7,2.8,3.3,79%,Workload,High
Product,3.9,3.5,3.4,3.3,3.8,4.1,90%,Resources,Low
Finance,3.7,3.4,3.5,3.6,3.2,3.8,85%,Automation tools,Low
HR,4.0,3.6,3.8,3.7,3.5,4.0,95%,Headcount,Low
Legal,3.5,3.3,3.4,3.1,3.0,3.6,78%,Workload,Medium
Operations,3.3,3.0,3.2,2.9,2.7,3.4,74%,Process improvement,High
IT,3.1,2.8,3.0,2.6,2.9,3.2,80%,Understaffing,High
""",
}

# --- Finance (4 files) ---

FILES["finance-budget-summary-2025.txt"] = {
    "tags": ["finance", "2025", "budget"],
    "dept": "finance",
    "content": """\
Finance Budget Summary - Fiscal Year 2025
===========================================

Prepared by: CFO David Park
Date: January 2025

Company-Wide Budget Overview
-----------------------------
Total FY2025 Budget: $180M
FY2024 Actual Spend: $164M
Increase: $16M (9.8%)

Revenue Context
---------------
FY2024 Revenue: $168M
FY2025 Revenue Forecast: $210M (see finance-revenue-forecast-2025.txt)
Budget-to-Revenue Ratio: 85.7% (target: below 85%)

The Q4 2024 revenue came in at $47M, exceeding the Q4 revenue target of
$42M (see sales-q4-report.txt). This overperformance provides a buffer
but does not eliminate the need for cost discipline.

Budget Allocation by Department
--------------------------------
Engineering:         $45M (25.0%)  - includes $12M requested in engineering-budget-proposal-2025.txt
Sales & Marketing:   $52M (28.9%)
Operations:          $23M (12.8%)
Customer Success:    $18M (10.0%)
G&A (Finance/HR/Legal/IT): $22M (12.2%)
Executive:            $5M (2.8%)
Contingency:         $15M (8.3%)

Budget Cuts Required
---------------------
To maintain the target budget-to-revenue ratio of 85%, we need to reduce
the proposed budget by $8M. Proposed cuts across departments:

1. Engineering: -$2M (reduce AI investment from $4.5M to $2.5M)
2. Sales & Marketing: -$2.5M (reduce marketing programs)
3. Operations: -$1M (defer automation projects)
4. G&A: -$1M (consolidate software licenses)
5. Contingency: -$1.5M (reduce from $15M to $13.5M)

Total proposed budget cuts: $8M

Impact Assessment
-----------------
The budget cuts carry significant risk:
- Engineering cut may delay AI features, impacting $15M in projected Q4 revenue
- Sales marketing cut may reduce pipeline generation by 10-15%
- See hr-hiring-plan-2025.txt for hiring plan impact scenarios

The board has requested final budget approval by February 15, 2025.
See executive-board-presentation-q4.txt for the board presentation.

Capital Expenditure
-------------------
Total CapEx: $12M
- Data center expansion: $5M
- Office renovations: $3M
- Equipment refresh: $4M
""",
}

FILES["finance-revenue-forecast-2025.txt"] = {
    "tags": ["finance", "2025", "revenue", "forecast"],
    "dept": "finance",
    "content": """\
Finance Revenue Forecast - Fiscal Year 2025
=============================================

Prepared by: FP&A Team
Date: January 2025

Annual Revenue Forecast: $210M
-------------------------------
This represents a 25% increase over FY2024 actual revenue of $168M.

Quarterly Revenue Targets
--------------------------
Q1 2025: $48M
  - Base: $44M from existing customers (renewal + expansion)
  - New: $4M from new customer acquisition
  - Dependency: Platform migration completion (engineering-roadmap-2025.txt)

Q2 2025: $50M
  - Base: $45M
  - New: $5M (API v3 drives partner channel revenue)
  - Dependency: API v3 launch on schedule

Q3 2025: $54M
  - Base: $47M
  - New: $7M (mobile redesign captures new segments)
  - Dependency: Mobile app launch

Q4 2025: $58M
  - Base: $48M
  - New: $10M (AI features drive enterprise upsell)
  - Dependency: AI integration launch (see engineering-roadmap-2025.txt)
  - Risk: If AI features are delayed, Q4 revenue target drops to $50M

The Q4 2025 revenue target of $58M is the most ambitious quarterly target
in company history, requiring both product execution and sales capacity.
The Q4 2024 revenue of $47M (see sales-q4-report.txt) gives confidence
but the $58M target is a 23% step-up.

Revenue by Segment Forecast
-----------------------------
Enterprise: $147M (70% of total, up from 64%)
Mid-Market: $42M (20%)
SMB: $21M (10%)

Revenue by Region Forecast
----------------------------
North America: $118M (56%)
Asia Pacific: $59M (28%)
Europe: $33M (16%)

Assumptions and Risks
----------------------
1. Assumes 120 new hires on schedule (see hr-hiring-plan-2025.txt)
2. Assumes no major product delays
3. Currency risk: 5% of APAC revenue denominated in local currencies
4. Competitive risk: Two competitors expected to launch similar products
5. Budget cuts could reduce marketing investment, lowering pipeline by 10-15%

Sensitivity Analysis
--------------------
Best case: $225M (+7% above forecast)
Base case: $210M (forecast)
Worst case: $185M (-12% below forecast, if AI delayed + hiring shortfall)
""",
}

FILES["finance-expense-report-q4.csv"] = {
    "tags": ["finance", "2024", "expenses"],
    "dept": "finance",
    "content": """\
Category,Department,Amount,ApprovedBy,Status,Notes
Cloud Infrastructure,Engineering,3200000,Sarah Chen,Paid,AWS and GCP Q4 usage
Sales Compensation,Sales,4650000,James Morrison,Paid,Base salary + Q4 commissions
Marketing Programs,Marketing,1800000,Lisa Zhang,Paid,Q4 campaign spend
Office Lease,Operations,1200000,COO,Paid,HQ and regional offices
Health Insurance,HR,2450000,Maria Santos,Paid,Q4 premium payments
Software Licenses,IT,890000,CTO,Paid,Annual renewals processed in Q4
Travel and Events,Sales,620000,James Morrison,Paid,Customer meetings and conferences
Legal Services,Legal,340000,General Counsel,Paid,Outside counsel fees
Recruitment,HR,580000,Maria Santos,Paid,Agency fees and job board postings
Equipment,IT,420000,CTO,Paid,Laptop refresh program
Training,HR,180000,Maria Santos,Paid,Professional development
Data Services,Engineering,290000,Sarah Chen,Paid,Third-party data APIs
Security Audit,Engineering,150000,Sarah Chen,Paid,Annual penetration testing
Contractor Payments,Engineering,380000,Sarah Chen,Paid,Specialized ML consultants
Customer Events,Customer Success,210000,VP CS,Paid,User conference Q4
Office Supplies,Operations,45000,COO,Paid,General office supplies
Insurance,Finance,320000,CFO,Paid,D&O and E&O insurance
Accounting,Finance,180000,CFO,Paid,External audit preparation
PR Agency,Marketing,120000,Lisa Zhang,Paid,Q4 media relations
Charity,Executive,50000,CEO,Paid,Annual charitable giving program
""",
}

FILES["finance-cash-flow-analysis.txt"] = {
    "tags": ["finance", "2025", "cash-flow"],
    "dept": "finance",
    "content": """\
Finance Cash Flow Analysis - FY2025
=====================================

Prepared by: Treasury Team
Date: January 2025

Current Cash Position
---------------------
Cash and equivalents: $42M (as of December 31, 2024)
Short-term investments: $15M
Total liquid assets: $57M

Monthly Burn Rate Analysis
---------------------------
Current monthly operating expenses: $13.7M
Projected monthly operating expenses (with hiring plan): $15.2M
Monthly revenue (projected average): $17.5M

Net monthly cash generation (projected): $2.3M

Cash Flow Projections by Quarter
---------------------------------
Q1 2025:
  Revenue: $48M
  Operating expenses: $44M
  Net cash flow: +$4M
  End-of-quarter cash: $61M

Q2 2025:
  Revenue: $50M
  Operating expenses: $46M (hiring ramp-up)
  Net cash flow: +$4M
  End-of-quarter cash: $65M

Q3 2025:
  Revenue: $54M
  Operating expenses: $49M (full hiring plan impact)
  Net cash flow: +$5M
  End-of-quarter cash: $70M

Q4 2025:
  Revenue: $58M
  Operating expenses: $48M (CapEx frontloaded earlier)
  Net cash flow: +$10M
  End-of-quarter cash: $80M

Risk: Cash Flow Crunch Scenario
--------------------------------
If the hiring plan (see hr-hiring-plan-2025.txt) is fully executed but
revenue comes in at the worst-case scenario of $185M (see
finance-revenue-forecast-2025.txt), a cash flow crunch could occur in Q3:

Worst-case Q3:
  Revenue: $43M (instead of $54M)
  Operating expenses: $49M
  Net cash flow: -$6M
  End-of-quarter cash: $49M

This would leave less than 4 months of runway, requiring either:
1. Emergency cost reductions (layoffs)
2. Debt financing
3. Accelerating the next funding round

Recommendation: Implement hiring milestones tied to revenue achievement.
If Q1 revenue falls below $44M, trigger a hiring pause for Q2.

Capital Expenditure Impact
--------------------------
CapEx of $12M (see finance-budget-summary-2025.txt) is planned for H1 2025.
This front-loads cash outflows but is necessary for data center expansion.
""",
}

# --- Legal (3 files) ---

FILES["legal-compliance-review.txt"] = {
    "tags": ["legal", "2025", "compliance"],
    "dept": "legal",
    "content": """\
Legal Compliance Review - FY2025
==================================

Prepared by: General Counsel Office
Date: January 2025

Regulatory Compliance Status
------------------------------

1. GDPR (General Data Protection Regulation)
   Status: PARTIALLY COMPLIANT
   Issues:
   - Data retention policies not fully implemented in 3 microservices
   - Cookie consent mechanism needs update for new EU guidance
   - Right-to-be-forgotten requests averaging 15 business days (target: 5)
   Risk: Fines up to 4% of global revenue ($8.4M based on FY2025 forecast)
   Action: Engineering team must implement data handling changes
   (see engineering-architecture-review.txt for technical requirements)

2. SOC 2 Type II
   Status: COMPLIANT
   Last audit: October 2024
   Next audit: October 2025
   No material findings in last audit.

3. CCPA (California Consumer Privacy Act)
   Status: COMPLIANT
   Data subject requests processed: 234 in 2024
   Average response time: 8 business days (within 45-day requirement)

4. ERISA / ACA (Benefits compliance)
   Status: COMPLIANT
   Note: Any benefits changes (see hr-benefits-review.txt) must be reviewed
   for ERISA compliance before implementation. The proposed dental plan
   downgrade requires 60-day employee notification.

5. PCI DSS (Payment Card Industry)
   Status: COMPLIANT
   Last assessment: December 2024
   Note: The payment service incident (see engineering-incident-report-q4.txt)
   did not result in any data breach, but follow-up assessment recommended.

6. Export Controls (ITAR/EAR)
   Status: COMPLIANT
   New markets in Southeast Asia and Middle East (see sales-strategy-2025.txt)
   require export control review before product launch.

Risk Summary
------------
Total regulatory risk exposure: approximately $15M in potential fines
Primary risk: GDPR non-compliance ($8.4M exposure)
Secondary risk: Data breach notification requirements

The GDPR risk has been added to the company risk register
(see risk-register-2025.txt, Risk #4).

Outside Counsel Budget
----------------------
FY2025 legal budget: $1.8M for outside counsel
- Regulatory advisory: $0.6M
- Employment law: $0.4M
- IP litigation defense: $0.5M
- M&A advisory: $0.3M
""",
}

FILES["legal-contract-summary-q4.txt"] = {
    "tags": ["legal", "2024", "contracts"],
    "dept": "legal",
    "content": """\
Legal Contract Summary - Q4 2024
==================================

Prepared by: Legal Department
Date: January 2025

Major Contracts Executed in Q4 2024
-------------------------------------
The following contracts were finalized in Q4 2024, contributing to the
Q4 revenue target achievement of $47M (see sales-q4-report.txt).

1. GlobalCorp Enterprise Agreement
   - Value: $4.2M over 3 years ($1.4M/year)
   - Type: Enterprise platform license
   - Key terms: Unlimited users, SLA 99.95%, dedicated support
   - Revenue recognition: Ratably over 36 months
   - Risk: SLA penalty clause of up to 10% per quarter for non-compliance

2. TechVentures Platform License
   - Value: $2.8M over 2 years ($1.4M/year)
   - Type: Platform + API access
   - Key terms: API rate limit of 10K requests/minute
   - Revenue recognition: Ratably over 24 months
   - Dependency: API v3 availability in Q2 2025

3. AsiaBank Digital Transformation
   - Value: $3.5M over 2 years
   - Type: Enterprise + professional services
   - Key terms: Data residency in Singapore, GDPR + local data protection
   - Revenue recognition: Milestone-based for services, ratable for license
   - Risk: Data residency requirement adds infrastructure cost

4. EuroManufacturing IoT Platform
   - Value: $1.9M over 1 year
   - Type: IoT platform license
   - Key terms: On-premise deployment option required
   - Revenue recognition: Upon delivery and acceptance
   - Risk: On-premise deployment not in current product roadmap

Total Q4 Contract Value: $12.4M
Contracts contributing to Q4 recognized revenue: $7.8M
Contracts contributing to future quarters: $4.6M

Contract Renewal Pipeline
--------------------------
$28M in contracts due for renewal in FY2025.
Renewal rate target: 92% (current: 89%)
At-risk renewals: $4.2M across 8 accounts (flagged to Customer Success)

IP Considerations
-----------------
Two contracts include IP indemnification clauses.
See legal-ip-portfolio-review.txt for current IP risk assessment.
""",
}

FILES["legal-ip-portfolio-review.txt"] = {
    "tags": ["legal", "2025", "ip"],
    "dept": "legal",
    "content": """\
Legal IP Portfolio Review - FY2025
=====================================

Prepared by: IP Counsel
Date: January 2025

Patent Portfolio
-----------------
Active patents: 23
Pending applications: 8
Provisional applications: 5
Total IP portfolio value (estimated): $45M

Key Patents
-----------
1. US Patent 11,234,567 - Real-time data processing pipeline
   Status: Active (expires 2038)
   Revenue attribution: Core to platform product

2. US Patent 11,345,678 - ML-based anomaly detection
   Status: Active (expires 2039)
   Revenue attribution: Enterprise security features

3. US Patent Application 17/456,789 - AI assistant interaction model
   Status: Pending (filed September 2024)
   Relevance: Critical for Q4 2025 AI feature launch

Infringement Risk Assessment
------------------------------
Risk Level: MEDIUM

1. Competitor B Patent Claim
   - Competitor B holds patents in mobile interaction patterns
   - Our mobile redesign (see engineering-roadmap-2025.txt Q3) may overlap
   - Risk: Patent infringement claim, potential damages $5-15M
   - Mitigation: Design-around analysis in progress, budget $200K

2. Open Source License Compliance
   - Audit identified 3 AGPL-licensed components in codebase
   - Risk: Potential obligation to open-source proprietary code
   - Mitigation: Replace with permissively-licensed alternatives

3. Trade Secret Protection
   - ML model training data and algorithms are trade secrets
   - Risk: Employee departure to competitor (see hr-employee-census-2025.txt
     for attrition data)
   - Mitigation: Updated non-compete and NDA agreements

These risks have been reflected in risk-register-2025.txt.

Trademark Portfolio
-------------------
Registered trademarks: 12
Pending applications: 3
No current opposition proceedings.
""",
}

# --- Executive / Cross-functional (5 files) ---

FILES["executive-board-presentation-q4.txt"] = {
    "tags": ["executive", "2024", "board"],
    "dept": "executive",
    "content": """\
Executive Board Presentation - Q4 2024 Results
=================================================

Prepared by: CEO Office
Date: January 2025

Agenda
------
1. Q4 2024 Financial Results
2. FY2024 Year in Review
3. FY2025 Strategic Plan
4. Budget Discussion
5. Risk Review

1. Q4 2024 Financial Results
------------------------------
Revenue: $47M vs $42M Q4 revenue target (+12% beat)
Operating Income: $8.2M (17.4% margin)
Net Income: $6.1M

Highlights:
- Record enterprise deal closings (see sales-q4-report.txt)
- Q4 revenue target exceeded despite November service outage
  (see engineering-incident-report-q4.txt, $2.3M revenue impact)
- Customer count grew to 2,340 (+18% YoY)

2. FY2024 Year in Review
--------------------------
Total Revenue: $168M (+23% YoY, exceeded $155M target)
Total Operating Income: $28M (16.7% margin)
Employee Count: 847 (see hr-employee-census-2025.txt)
Customer NPS: 52 (industry average: 44)

3. FY2025 Strategic Plan
--------------------------
Revenue Target: $210M (+25% YoY)
See finance-revenue-forecast-2025.txt for quarterly breakdown.
See sales-strategy-2025.txt for go-to-market plan.
See engineering-roadmap-2025.txt for product roadmap.

Three strategic pillars:
a) Growth: Enterprise expansion, geographic markets, partner channel
b) Efficiency: Platform modernization, process automation
c) Innovation: AI integration, mobile redesign

4. Budget Discussion
---------------------
Proposed Budget: $180M
Budget Cuts Required: $8M to maintain target ratios
See finance-budget-summary-2025.txt for detailed breakdown.

Key debate points:
- Engineering AI investment ($4.5M requested, $2.5M proposed after cuts)
- Risk: Budget cuts to AI may reduce Q4 2025 revenue by $15M
- Hiring plan: 120 new hires planned (see hr-hiring-plan-2025.txt)
- Benefits changes proposed to save $750K (see hr-benefits-review.txt)

5. Risk Review
--------------
Top 5 risks for FY2025 (see risk-register-2025.txt for full list):
1. Cybersecurity breach / data loss
2. Key talent attrition (especially ML/AI engineers)
3. Revenue target miss if AI features delayed
4. GDPR compliance gap (see legal-compliance-review.txt)
5. Competitive pressure from well-funded rivals

Board Action Items:
- Approve FY2025 budget by February 15
- Review hiring plan milestones
- Schedule cybersecurity review for Q1 board meeting
""",
}

FILES["risk-register-2025.txt"] = {
    "tags": ["executive", "2025", "risk"],
    "dept": "executive",
    "content": """\
Company Risk Register - FY2025
================================

Prepared by: Risk Management Committee
Date: January 2025

Risk Rating Scale: Impact (1-5) x Likelihood (1-5) = Risk Score (1-25)

Risk #1: Cybersecurity Breach
  Impact: 5 | Likelihood: 3 | Score: 15
  Description: Data breach exposing customer PII or financial data
  Current controls: SOC 2 compliance, annual pen testing, WAF, encryption
  Additional mitigation: Increase security team headcount by 2
  Owner: CISO
  Reference: engineering-architecture-review.txt (data handling)

Risk #2: Key Person Dependency
  Impact: 4 | Likelihood: 4 | Score: 16
  Description: Loss of critical technical talent, especially ML/AI engineers
  Current controls: Retention bonuses, equity refresh, career development
  Additional mitigation: Cross-training program, documentation initiative
  Owner: CHRO
  Reference: hr-employee-census-2025.txt (22% ML/AI attrition rate)

Risk #3: Market Downturn
  Impact: 4 | Likelihood: 2 | Score: 8
  Description: Economic recession reducing enterprise IT spending
  Current controls: Diversified customer base, multi-year contracts
  Additional mitigation: Accelerate SMB product-led growth
  Owner: CFO
  Reference: finance-revenue-forecast-2025.txt (sensitivity analysis)

Risk #4: Regulatory Non-Compliance (GDPR)
  Impact: 5 | Likelihood: 3 | Score: 15
  Description: GDPR fine of up to 4% of global revenue
  Current controls: DPO appointed, privacy impact assessments
  Additional mitigation: Complete data handling remediation by Q2
  Owner: General Counsel
  Reference: legal-compliance-review.txt

Risk #5: Supply Chain Disruption (Cloud)
  Impact: 3 | Likelihood: 2 | Score: 6
  Description: Cloud provider outage or pricing increase
  Current controls: Multi-cloud architecture, reserved instances
  Additional mitigation: Evaluate private cloud for critical workloads
  Owner: CTO

Risk #6: Talent Retention
  Impact: 4 | Likelihood: 3 | Score: 12
  Description: High attrition rate impacting productivity and hiring costs
  Current controls: Competitive compensation, benefits package
  Additional mitigation: Address satisfaction survey concerns
  Owner: CHRO
  Reference: hr-employee-satisfaction-survey.csv

Risk #7: Budget Overrun
  Impact: 3 | Likelihood: 3 | Score: 9
  Description: Departments exceeding allocated budgets
  Current controls: Monthly budget reviews, approval workflows
  Additional mitigation: Implement automated budget alerts
  Owner: CFO
  Reference: finance-budget-summary-2025.txt

Risk #8: Technology Debt
  Impact: 3 | Likelihood: 4 | Score: 12
  Description: Legacy systems causing reliability issues and slowing development
  Current controls: Architecture review board, tech debt sprints
  Additional mitigation: Platform migration in Q1
  Owner: CTO
  Reference: engineering-architecture-review.txt, engineering-incident-report-q4.txt

Risk #9: Competitive Pressure
  Impact: 4 | Likelihood: 4 | Score: 16
  Description: Competitors launching AI features, acquiring companies
  Current controls: Product differentiation, enterprise relationships
  Additional mitigation: Accelerate AI integration (Q4 2025)
  Owner: CEO
  Reference: sales-strategy-2025.txt (competitive landscape)

Risk #10: Currency Fluctuation
  Impact: 2 | Likelihood: 3 | Score: 6
  Description: APAC revenue denominated in local currencies
  Current controls: Natural hedging through local expenses
  Additional mitigation: Evaluate currency hedging instruments
  Owner: CFO
  Reference: finance-revenue-forecast-2025.txt

Summary: Top risks by score
  1. Key Person Dependency: 16
  2. Competitive Pressure: 16
  3. Cybersecurity Breach: 15
  4. Regulatory Non-Compliance: 15
  5. Technology Debt: 12
  6. Talent Retention: 12
""",
}

FILES["strategic-priorities-2025.txt"] = {
    "tags": ["executive", "2025", "strategy"],
    "dept": "executive",
    "content": """\
Strategic Priorities - Fiscal Year 2025
=========================================

Prepared by: CEO
Date: January 2025

Our Three Pillars for 2025
---------------------------

PILLAR 1: GROWTH
  Target: $210M revenue (+25% YoY)
  Key initiatives:
  - Enterprise expansion (grow from 64% to 70% of revenue)
  - Geographic expansion: Southeast Asia, Middle East, Latin America
  - Partner channel: Target 15% of new bookings from partners
  - Pricing: Introduce consumption-based pricing for API products

  Budget allocation: $52M (Sales & Marketing)
  Headcount: 35 new sales hires, 20 new marketing hires
  See: sales-strategy-2025.txt, hr-hiring-plan-2025.txt

PILLAR 2: EFFICIENCY
  Target: Improve operating margin from 16.7% to 18%
  Key initiatives:
  - Platform migration to reduce infrastructure costs by 20%
  - Process automation across all departments
  - Vendor consolidation (target: reduce software spend by 15%)
  - Budget discipline: $8M in planned budget cuts

  Budget allocation: Cross-departmental
  See: finance-budget-summary-2025.txt, engineering-roadmap-2025.txt

PILLAR 3: INNOVATION
  Target: Launch 3 major product capabilities
  Key initiatives:
  - API v3 with GraphQL (Q2)
  - Mobile redesign (Q3)
  - AI assistant feature (Q4)
  - File of new patents (target: 10 new applications)

  Budget allocation: $12M (Engineering)
  Headcount: 45 new engineering hires
  See: engineering-roadmap-2025.txt, engineering-budget-proposal-2025.txt

Success Metrics
---------------
1. Revenue: $210M
2. Operating margin: 18%
3. Employee count: ~967 (847 + 120 hires)
4. Customer NPS: 55 (up from 52)
5. Enterprise customer count: +30% YoY

Risks to Strategy
-----------------
The strategy depends on three critical assumptions:
1. Engineering delivers AI features on schedule for Q4
2. Hiring plan is executed without major delays
3. Budget cuts do not impair revenue-generating activities

See risk-register-2025.txt for detailed risk assessment.
""",
}

FILES["all-hands-meeting-notes-jan2025.txt"] = {
    "tags": ["executive", "2025", "meeting"],
    "dept": "executive",
    "content": """\
All-Hands Meeting Notes - January 15, 2025
=============================================

CEO Opening Remarks
-------------------
"2024 was a breakthrough year for us. We exceeded our revenue target of
$155M, finishing the year at $168M. Q4 was particularly strong — we hit
$47M against a Q4 revenue target of $42M. Thank you all for the incredible
work that made this possible."

FY2024 Highlights
-----------------
- Revenue: $168M (+23% YoY)
- Customer count: 2,340 (+18% YoY)
- Employee count: 847
- Key wins: GlobalCorp ($4.2M), AsiaBank ($3.5M)
- Product: Launched 47 features, 99.95% uptime

Addressing the November Outage
-------------------------------
"I want to acknowledge the November 15th outage directly. We lost $2.3M
in transactions, and more importantly, we let our customers down. Sarah
and the engineering team have already implemented fixes, and we're investing
in resilience to make sure this doesn't happen again."
(See engineering-incident-report-q4.txt for details)

FY2025 Plans
------------
"Our target for 2025 is $210M in revenue. That's ambitious — 25% growth.
Here's how we'll get there:"

1. Product roadmap: Platform migration (Q1), API v3 (Q2), mobile (Q3), AI (Q4)
   (see engineering-roadmap-2025.txt)

2. Hiring: We plan to hire 120 new team members this year
   (see hr-hiring-plan-2025.txt)

3. Budget: Our total budget is $180M, and yes, we're implementing some
   budget cuts totaling $8M to maintain financial discipline
   (see finance-budget-summary-2025.txt)

Budget Cuts Rationale
---------------------
"I know budget cuts are never popular. Let me be transparent about why:
Our budget-to-revenue ratio needs to stay below 85%. The cuts are spread
across departments and are designed to minimize impact on revenue-generating
activities. Engineering will see a $2M reduction, which primarily affects
the pace of AI investment. Sales and marketing will see a $2.5M reduction
in campaign budgets."

Employee Q&A Highlights
------------------------
Q: "Will there be layoffs?"
A: "No. We are hiring 120 people. The budget cuts are about prioritization,
   not reducing headcount."

Q: "Will benefits change?"
A: "We are reviewing some benefits to find savings. The dental plan may
   move to a standard tier, saving about $420K. We'll communicate any
   changes with 60 days notice."
   (See hr-benefits-review.txt)

Q: "How serious is the competitive threat from Competitor A's AI launch?"
A: "Very serious, which is why even with budget cuts, we're still investing
   $2.5M in AI for 2025. Our Q4 AI launch is critical."
   (See risk-register-2025.txt, Risk #9)

Closing
-------
"We have an incredible opportunity ahead of us. $210M is achievable if
we execute well. Let's make 2025 our best year yet."
""",
}

FILES["cross-department-project-tracker.csv"] = {
    "tags": ["executive", "2025", "projects"],
    "dept": "executive",
    "content": """\
Project,Lead,Department,Status,Budget,Deadline,Priority,Dependencies,RiskLevel
Platform Migration,Sarah Chen,Engineering,In Progress,3200000,2025-03-31,Critical,Cloud infrastructure,High
API v3 Launch,API Team Lead,Engineering,Planning,1800000,2025-06-30,Critical,Platform migration,Medium
Mobile Redesign,Mobile Lead,Engineering,Planning,2500000,2025-09-30,High,API v3,Medium
AI Integration,ML Team Lead,Engineering,Research,2500000,2025-12-31,High,ML infrastructure,High
Enterprise Sales Program,James Morrison,Sales,Active,8000000,2025-12-31,Critical,Product features,Medium
Partner Channel Launch,Partner Lead,Sales,Planning,1500000,2025-06-30,High,API v3,Medium
Geographic Expansion,Regional Lead,Sales,Planning,2000000,2025-09-30,Medium,Legal review,Medium
Brand Refresh,Lisa Zhang,Marketing,Active,800000,2025-03-31,Medium,None,Low
GDPR Remediation,Legal Counsel,Legal,In Progress,600000,2025-06-30,Critical,Engineering support,High
Benefits Optimization,HR Team,HR,Planning,100000,2025-04-30,Medium,Legal review,Low
Vendor Consolidation,IT Director,IT,Planning,200000,2025-06-30,Medium,None,Low
Office Renovation,Facilities,Operations,Not Started,3000000,2025-09-30,Low,Budget approval,Low
""",
}


# ============================================================
# Format-aware file creation
# ============================================================


def _fix_cross_references(content: str) -> str:
    """Replace .txt references with correct format extensions."""
    for base, fmt in FORMAT_MAP.items():
        content = content.replace(f"{base}.txt", f"{base}.{fmt}")
    return content


def _parse_rst_lines(content: str) -> None:
    """Parse RST-style text into (type, text) tuples.
    Types: 'h1', 'h2', 'text', 'blank'."""
    lines = content.split("\n")
    result = []
    i = 0
    while i < len(lines):
        line = lines[i]
        next_line = lines[i + 1] if i + 1 < len(lines) else ""
        stripped_next = next_line.strip()
        if stripped_next and len(stripped_next) >= 3 and set(stripped_next) <= {"="}:
            result.append(("h1", line.strip()))
            i += 2
            continue
        if stripped_next and len(stripped_next) >= 3 and set(stripped_next) <= {"-"}:
            result.append(("h2", line.strip()))
            i += 2
            continue
        if line.strip():
            result.append(("text", line))
        else:
            result.append(("blank", ""))
        i += 1
    return result


def _create_docx(filepath: str, content: str) -> None:
    """Create a .docx file from text content."""
    from docx import Document

    doc = Document()
    for kind, text in _parse_rst_lines(content):
        if kind == "h1":
            doc.add_heading(text, level=1)
        elif kind == "h2":
            doc.add_heading(text, level=2)
        elif kind == "text":
            doc.add_paragraph(text)
        # skip blanks
    doc.save(filepath)


def _create_pdf(filepath: str, content: str) -> None:
    """Create a .pdf file from text content using PyMuPDF."""
    import fitz

    doc = fitz.open()
    page = doc.new_page(width=595.28, height=841.89)  # A4
    y = 72

    for kind, text in _parse_rst_lines(content):
        if kind == "h1":
            if y > 750:
                page = doc.new_page(width=595.28, height=841.89)
                y = 72
            page.insert_text(fitz.Point(72, y), text, fontsize=16, fontname="hebo")
            y += 28
        elif kind == "h2":
            if y > 760:
                page = doc.new_page(width=595.28, height=841.89)
                y = 72
            page.insert_text(fitz.Point(72, y), text, fontsize=13, fontname="hebo")
            y += 22
        elif kind == "text":
            if y > 780:
                page = doc.new_page(width=595.28, height=841.89)
                y = 72
            page.insert_text(fitz.Point(72, y), text, fontsize=10, fontname="helv")
            y += 14
        elif kind == "blank":
            y += 8

    doc.save(filepath)
    doc.close()


def _create_pptx(filepath: str, content: str) -> None:
    """Create a .pptx file from text content."""
    from pptx import Presentation
    from pptx.util import Pt

    prs = Presentation()
    sections = []
    current = {"title": "", "lines": []}

    for kind, text in _parse_rst_lines(content):
        if kind in ("h1", "h2"):
            if current["title"] or current["lines"]:
                sections.append(current)
            current = {"title": text, "lines": []}
        elif kind == "text":
            current["lines"].append(text)

    if current["title"] or current["lines"]:
        sections.append(current)

    for section in sections:
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content
        slide.shapes.title.text = section["title"] or "Overview"
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.clear()
        for j, line in enumerate(section["lines"][:20]):
            if j == 0:
                tf.text = line.strip()
            else:
                p = tf.add_paragraph()
                p.text = line.strip()
                p.font.size = Pt(14)

    prs.save(filepath)


# ============================================================
# File generation and upload
# ============================================================


def generate_files() -> list[dict]:
    """Generate all workspace files in proper office formats."""
    os.makedirs(WORKSPACE_DIR, exist_ok=True)
    result = []

    for old_filename, info in FILES.items():
        base = os.path.splitext(old_filename)[0]
        fmt = FORMAT_MAP.get(base, os.path.splitext(old_filename)[1].lstrip("."))
        new_filename = f"{base}.{fmt}"
        filepath = os.path.join(WORKSPACE_DIR, new_filename)

        content = _fix_cross_references(info["content"])

        if fmt == "csv":
            with open(filepath, "w", encoding="utf-8") as f:
                f.write(content)
        elif fmt == "docx":
            _create_docx(filepath, content)
        elif fmt == "pdf":
            _create_pdf(filepath, content)
        elif fmt == "pptx":
            _create_pptx(filepath, content)

        result.append({
            "path": filepath,
            "filename": new_filename,
            "tags": info["tags"],
            "dept": info["dept"],
        })

    return result


async def upload_to_filedb(files: list[dict], base_url: str) -> None:
    """Upload all generated files to FileDB."""
    async with httpx.AsyncClient(base_url=base_url, timeout=60) as client:
        # Health check
        try:
            resp = await client.get("/health")
            if resp.status_code != 200:
                print(f"  WARNING: /health returned {resp.status_code}")
        except httpx.ConnectError:
            print(f"  ERROR: Cannot connect to {base_url}")
            print("  Make sure FileDB is running.")
            sys.exit(1)

        for f in files:
            with open(f["path"], "rb") as fp:
                resp = await client.post(
                    "/files",
                    files={"file": (f["filename"], fp)},
                    data={
                        "tags": json.dumps(f["tags"]),
                        "metadata": json.dumps({"department": f["dept"]}),
                    },
                )
            if resp.status_code in (200, 201):
                data = resp.json()
                status = data.get("status", "?")
                print(f"  {f['filename']}: {status}")
            else:
                print(f"  {f['filename']}: ERROR {resp.status_code} - {resp.text[:100]}")


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Generate FileDB benchmark workspace")
    parser.add_argument("--url", default=FILEDB_URL, help="FileDB base URL")
    parser.add_argument("--skip-upload", action="store_true", help="Only generate files, don't upload")
    args = parser.parse_args()

    print("Generating workspace files...")
    files = generate_files()
    print(f"Generated {len(files)} files in {WORKSPACE_DIR}")

    if not args.skip_upload:
        print(f"\nUploading to FileDB at {args.url}...")
        asyncio.run(upload_to_filedb(files, args.url))
        print("Done.")
    else:
        print("Skipping upload (--skip-upload).")
