"""PPTX parser using python-pptx.

One page per slide. Extracts title, body text, speaker notes, tables.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from opendb_core.parsers.base import Page, ParseResult
from opendb_core.parsers.registry import register


class PptxParser:
    def parse(self, file_path: Path) -> ParseResult:
        prs = Presentation(str(file_path))
        pages: list[Page] = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            title = self._get_slide_title(slide)
            body_parts: list[str] = []

            # Extract text from all shapes
            for shape in slide.shapes:
                if shape.has_text_frame:
                    # Skip title shape (already captured)
                    if shape == slide.shapes.title:
                        continue
                    text = shape.text_frame.text.strip()
                    if text:
                        body_parts.append(text)

                if shape.has_table:
                    table_text = self._format_table(shape.table)
                    if table_text:
                        body_parts.append(table_text)

            # Speaker notes
            notes = self._get_notes(slide)
            if notes:
                body_parts.append(f"[Speaker Notes]\n{notes}")

            text = "\n\n".join(body_parts)
            if title:
                text = f"{title}\n\n{text}" if text else title

            pages.append(
                Page(
                    page_number=slide_num,
                    section_title=title,
                    text=text.strip(),
                )
            )

        metadata = self._extract_metadata(prs)
        return ParseResult(pages=pages, extracted_metadata=metadata)

    def _get_slide_title(self, slide) -> str | None:
        """Get slide title from the title placeholder."""
        if slide.shapes.title and slide.shapes.title.text.strip():
            return slide.shapes.title.text.strip()
        return None

    def _get_notes(self, slide) -> str:
        """Extract speaker notes from a slide."""
        try:
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                return notes_slide.notes_text_frame.text.strip()
        except (AttributeError, RuntimeError):
            pass
        return ""

    def _format_table(self, table) -> str:
        """Format a PPTX table as pipe-delimited text."""
        rows: list[str] = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append("| " + " | ".join(cells) + " |")

        if len(rows) > 1:
            header_cells = list(table.rows[0].cells)
            separator = "|" + "|".join(
                "-" * (len(c.text.strip()) + 2) for c in header_cells
            ) + "|"
            rows.insert(1, separator)

        return "\n".join(rows)

    def _extract_metadata(self, prs: Presentation) -> dict:
        """Extract PPTX core properties."""
        props = prs.core_properties
        result = {}

        if props.author:
            result["author"] = props.author
        if props.title:
            result["title"] = props.title
        if hasattr(props, "company") and getattr(props, "company", None):
            result["company"] = props.company
        if props.created:
            result["created"] = props.created.isoformat()
        if props.modified:
            result["modified"] = props.modified.isoformat()

        result["slide_count"] = len(prs.slides)
        return result


register(
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    PptxParser(),
)
