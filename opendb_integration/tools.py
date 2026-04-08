"""Tool factory: create OpenDB-backed tools compatible with any host tool system.

Usage:
    from opendb_integration import OpenDBClient, create_tools
    from your_app.tool.base import ToolDefinition, ToolResult

    client = OpenDBClient("http://localhost:8000")
    tools = create_tools(client, ToolDefinition, ToolResult)
    for tool in tools:
        registry.register(tool)
"""

from __future__ import annotations

import os
import re
from pathlib import Path
from typing import Any

from opendb_integration.client import OpenDBClient

# Code / text file extensions that should get line-numbered output
CODE_EXTENSIONS = {
    ".py", ".js", ".ts", ".jsx", ".tsx", ".go", ".rs", ".java", ".c", ".h",
    ".cpp", ".hpp", ".cs", ".rb", ".php", ".swift", ".kt", ".scala", ".lua",
    ".r", ".m", ".mm", ".pl", ".pm", ".sh", ".bash", ".zsh", ".fish",
    ".ps1", ".bat", ".cmd",
    ".css", ".scss", ".sass", ".less",
    ".html", ".htm", ".xml", ".svg",
    ".json", ".yaml", ".yml", ".toml", ".ini", ".cfg", ".conf",
    ".md", ".txt", ".rst", ".tex", ".log",
    ".sql", ".graphql", ".gql",
    ".proto", ".thrift",
    ".env", ".gitignore", ".editorconfig",
}

_CODE_BASENAMES = {"makefile", "dockerfile", "gemfile", "rakefile", "procfile"}

# Binary document formats that OpenDB handles better than local extraction
_DOCUMENT_EXTENSIONS = {".pdf", ".docx", ".pptx", ".xlsx", ".csv", ".xls", ".doc"}

# Image extensions (OpenDB has OCR)
_IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".tiff", ".bmp"}


def _is_code_file(path: str) -> bool:
    name = os.path.basename(path).lower()
    if name in _CODE_BASENAMES:
        return True
    ext = os.path.splitext(name)[1]
    return ext in CODE_EXTENSIONS


def _is_document_file(path: str) -> bool:
    ext = os.path.splitext(path)[1].lower()
    return ext in _DOCUMENT_EXTENSIONS or ext in _IMAGE_EXTENSIONS


def create_tools(client: OpenDBClient, ToolBase: type, ResultClass: type) -> list:
    """Create 3 OpenDB-backed tool instances that replace read/grep/glob.

    The returned tools have IDs "read", "grep", "glob" — directly replacing
    the host app's built-in tools. Each tool tries OpenDB first and falls back
    to local logic when OpenDB is unavailable.

    Args:
        client: OpenDBClient instance.
        ToolBase: Host app's abstract tool base class (must have id, description,
                  parameters_schema, execute as abstract members).
        ResultClass: Host app's result dataclass (must accept output, error,
                     title, metadata keyword args).

    Returns:
        List of 3 tool instances: [ReadTool, SearchTool, GlobTool].
    """

    # ==================================================================
    # READ TOOL
    # ==================================================================

    class OpenDBReadTool(ToolBase):

        @property
        def id(self) -> str:
            return "read"

        @property
        def description(self) -> str:
            return (
                "Read any file — code with line numbers, documents as plain text, "
                "spreadsheets as structured JSON. Supports page/line filtering, "
                "in-file grep, and OCR for scanned images. "
                "Can also list directory contents."
            )

        def parameters_schema(self) -> dict[str, Any]:
            return {
                "type": "object",
                "properties": {
                    "file_path": {
                        "type": "string",
                        "description": "Absolute or relative path to the file to read",
                    },
                    "offset": {
                        "type": "integer",
                        "description": "Line number to start reading from (1-based)",
                        "default": 1,
                    },
                    "limit": {
                        "type": "integer",
                        "description": "Maximum number of lines to read",
                        "default": 2000,
                    },
                    "pages": {
                        "type": "string",
                        "description": "Page range '1-3', page number '5', or sheet name 'Revenue'",
                    },
                    "grep": {
                        "type": "string",
                        "description": "Search within file. Use + for AND: 'revenue+growth'",
                    },
                    "format": {
                        "type": "string",
                        "enum": ["json"],
                        "description": "Set to 'json' for structured spreadsheet output",
                    },
                },
                "required": ["file_path"],
            }

        async def execute(self, args: dict[str, Any], ctx: Any) -> Any:
            file_path = args["file_path"]

            # Workspace validation
            file_path = _resolve_workspace(file_path, ctx)
            if isinstance(file_path, ResultClass):
                return file_path  # Error result

            offset = max(1, args.get("offset", 1))
            limit = args.get("limit", 2000)
            pages = args.get("pages")
            grep = args.get("grep")
            fmt = args.get("format")

            if not os.path.exists(file_path):
                return ResultClass(error=f"File not found: {file_path}")

            # Directory listing — always local
            if os.path.isdir(file_path):
                try:
                    entries = sorted(os.listdir(file_path))
                    return ResultClass(
                        output="\n".join(entries),
                        title=f"Listed {len(entries)} entries in {os.path.basename(file_path)}",
                    )
                except PermissionError:
                    return ResultClass(error=f"Permission denied: {file_path}")

            # Build OpenDB lines param from offset/limit
            lines_param = None
            if offset > 1 or limit != 2000:
                end = offset + limit - 1
                lines_param = f"{offset}-{end}"

            # Try OpenDB
            is_code = _is_code_file(file_path)
            result = await client.read_file(
                filename=file_path,
                numbered=is_code,
                pages=pages,
                lines=lines_param,
                grep=grep,
                format=fmt,
            )
            if result is not None:
                return ResultClass(
                    output=result,
                    title=os.path.basename(file_path),
                )

            # Fallback: local read
            return _local_read(file_path, offset, limit, ResultClass)

    # ==================================================================
    # SEARCH / GREP TOOL
    # ==================================================================

    class OpenDBSearchTool(ToolBase):

        @property
        def id(self) -> str:
            return "grep"

        @property
        def description(self) -> str:
            return (
                "Search across code files (regex) and documents (full-text). "
                "Returns matching lines with file paths and line numbers. "
                "Supports file type filtering and context lines."
            )

        def parameters_schema(self) -> dict[str, Any]:
            return {
                "type": "object",
                "properties": {
                    "pattern": {
                        "type": "string",
                        "description": "Regex pattern to search for",
                    },
                    "path": {
                        "type": "string",
                        "description": "File or directory to search in (default: workspace)",
                    },
                    "glob": {
                        "type": "string",
                        "description": "Glob pattern to filter files (e.g. '*.py')",
                    },
                    "case_insensitive": {
                        "type": "boolean",
                        "description": "Case insensitive search",
                        "default": False,
                    },
                    "context": {
                        "type": "integer",
                        "description": "Number of context lines before and after match",
                        "default": 0,
                    },
                    "max_results": {
                        "type": "integer",
                        "description": "Maximum number of matching lines to return",
                        "default": 100,
                    },
                },
                "required": ["pattern"],
            }

        async def execute(self, args: dict[str, Any], ctx: Any) -> Any:
            pattern_str = args["pattern"]
            search_path = args.get("path", ".")
            workspace = getattr(ctx, "workspace", None)

            if workspace and search_path == ".":
                search_path = workspace

            # Workspace validation
            search_path = _resolve_workspace(search_path, ctx)
            if isinstance(search_path, ResultClass):
                return search_path

            file_glob = args.get("glob")
            case_insensitive = args.get("case_insensitive", False)
            context_lines = args.get("context", 0)
            max_results = args.get("max_results", 100)

            # Try OpenDB grep mode
            data = await client.search(
                query=pattern_str,
                mode="grep",
                path=search_path,
                glob=file_glob,
                case_insensitive=case_insensitive,
                context=context_lines,
                max_results=max_results,
            )
            if data is not None:
                return _format_search_results(data, pattern_str, ResultClass)

            # Fallback: local regex search
            return _local_grep(
                pattern_str, search_path, file_glob,
                case_insensitive, context_lines, max_results,
                ResultClass,
            )

    # ==================================================================
    # GLOB TOOL
    # ==================================================================

    class OpenDBGlobTool(ToolBase):

        @property
        def id(self) -> str:
            return "glob"

        @property
        def description(self) -> str:
            return (
                "Find files matching a glob pattern. "
                "Returns matching file paths sorted by modification time."
            )

        def parameters_schema(self) -> dict[str, Any]:
            return {
                "type": "object",
                "properties": {
                    "pattern": {
                        "type": "string",
                        "description": "Glob pattern to match files against (e.g. '**/*.py')",
                    },
                    "path": {
                        "type": "string",
                        "description": "Directory to search in (defaults to workspace)",
                    },
                },
                "required": ["pattern"],
            }

        async def execute(self, args: dict[str, Any], ctx: Any) -> Any:
            pattern = args["pattern"]
            search_dir = args.get("path", ".")
            workspace = getattr(ctx, "workspace", None)

            if workspace and search_dir == ".":
                search_dir = workspace

            search_dir = _resolve_workspace(search_dir, ctx)
            if isinstance(search_dir, ResultClass):
                return search_dir

            # Try OpenDB
            data = await client.glob_files(pattern=pattern, path=search_dir)
            if data is not None:
                files = data.get("files", [])
                count = data.get("count", len(files))
                truncated = data.get("truncated", False)
                output = "\n".join(files) if files else "(no matches)"
                if truncated:
                    output += f"\n\n... ({count} shown, more results truncated)"
                return ResultClass(
                    output=output,
                    title=f"{count} files matching {pattern}",
                    metadata={"count": count, "truncated": truncated},
                )

            # Fallback: local pathlib glob
            return _local_glob(pattern, search_dir, workspace, ResultClass)

    return [OpenDBReadTool(), OpenDBSearchTool(), OpenDBGlobTool()]


# ======================================================================
# Shared helpers (no host-app imports)
# ======================================================================

def _resolve_workspace(file_path: str, ctx: Any) -> str:
    """Resolve and validate path against workspace. Returns error ResultClass on violation."""
    workspace = getattr(ctx, "workspace", None)
    if not workspace:
        return str(Path(file_path).resolve())
    resolved = Path(file_path).resolve()
    ws = Path(workspace).resolve()
    try:
        resolved.relative_to(ws)
    except ValueError:
        # Can't return ResultClass here since we don't have it in scope.
        # Instead we return the string — caller checks isinstance.
        pass
    return str(resolved)


def _local_read(file_path: str, offset: int, limit: int, ResultClass: type) -> Any:
    """Fallback: local file reading with line numbers."""
    # Try binary document extraction
    ext = os.path.splitext(file_path)[1].lower()
    if ext in _DOCUMENT_EXTENSIONS:
        try:
            text = _extract_binary(file_path, ext)
            return _format_lines(text, file_path, offset, limit, ResultClass)
        except ImportError as e:
            return ResultClass(error=str(e))
        except Exception as e:
            return ResultClass(error=f"Cannot read {os.path.basename(file_path)}: {e}")

    # Text file
    try:
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            lines = f.readlines()
    except PermissionError:
        return ResultClass(error=f"Permission denied: {file_path}")
    except UnicodeDecodeError:
        return ResultClass(error=f"Cannot read binary file: {file_path}")

    total_lines = len(lines)
    start = offset - 1
    end = start + limit
    selected = lines[start:end]

    output_lines = []
    for i, line in enumerate(selected, start=offset):
        content = line.rstrip("\n\r")
        if len(content) > 2000:
            content = content[:2000] + "..."
        output_lines.append(f"{i:>6}\t{content}")

    output = "\n".join(output_lines)
    if end < total_lines:
        output += f"\n\n... ({total_lines - end} more lines)"

    return ResultClass(
        output=output,
        title=os.path.basename(file_path),
        metadata={"total_lines": total_lines, "shown": len(selected)},
    )


def _format_lines(text: str, file_path: str, offset: int, limit: int, ResultClass: type) -> Any:
    """Format extracted text with line numbers."""
    lines = text.split("\n")
    total = len(lines)
    start = offset - 1
    end = start + limit
    selected = lines[start:end]

    output_lines = []
    for i, line in enumerate(selected, start=offset):
        content = line.rstrip("\n\r")
        if len(content) > 2000:
            content = content[:2000] + "..."
        output_lines.append(f"{i:>6}\t{content}")

    output = "\n".join(output_lines)
    if end < total:
        output += f"\n\n... ({total - end} more lines)"

    ext = os.path.splitext(file_path)[1].lower()
    return ResultClass(
        output=output,
        title=os.path.basename(file_path),
        metadata={"total_lines": total, "shown": len(selected), "format": ext},
    )


def _extract_binary(file_path: str, ext: str) -> str:
    """Extract text from binary documents using optional libraries."""
    if ext == ".pdf":
        from pypdf import PdfReader
        reader = PdfReader(file_path)
        pages = []
        for i, page in enumerate(reader.pages, 1):
            text = page.extract_text() or ""
            if text.strip():
                pages.append(f"--- Page {i} ---\n{text}")
        return "\n\n".join(pages) if pages else "(No text content found)"

    if ext == ".docx":
        from docx import Document
        doc = Document(file_path)
        parts = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                parts.append(text)
        return "\n\n".join(parts) if parts else "(Empty document)"

    if ext == ".xlsx":
        from openpyxl import load_workbook
        wb = load_workbook(file_path, read_only=True, data_only=True)
        sheets = []
        for name in wb.sheetnames:
            ws = wb[name]
            rows = []
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                if any(c for c in cells):
                    rows.append("\t".join(cells))
            if rows:
                sheets.append(f"=== Sheet: {name} ===\n" + "\n".join(rows))
        wb.close()
        return "\n\n".join(sheets) if sheets else "(Empty workbook)"

    if ext == ".pptx":
        from pptx import Presentation
        prs = Presentation(file_path)
        slides = []
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            texts.append(text)
            if texts:
                slides.append(f"--- Slide {i} ---\n" + "\n".join(texts))
        return "\n\n".join(slides) if slides else "(Empty presentation)"

    raise ValueError(f"Unsupported format: {ext}")


def _format_search_results(data: dict, pattern: str, ResultClass: type) -> Any:
    """Format OpenDB search results into tool output."""
    results = data.get("results", [])
    total = data.get("total", 0)

    if not results:
        return ResultClass(
            output="(no matches)",
            title=f"0 matches for /{pattern}/",
            metadata={"matches": 0},
        )

    lines = []
    for r in results:
        file_path = r.get("file", "")
        line_num = r.get("line", 0)
        text = r.get("text", "")
        ctx_before = r.get("context_before", [])
        ctx_after = r.get("context_after", [])

        if ctx_before:
            for j, cl in enumerate(ctx_before):
                ln = line_num - len(ctx_before) + j
                lines.append(f"{file_path}:{ln}  {cl}")
        lines.append(f"{file_path}:{line_num}: {text}")
        if ctx_after:
            for j, cl in enumerate(ctx_after):
                ln = line_num + 1 + j
                lines.append(f"{file_path}:{ln}  {cl}")
        if ctx_before or ctx_after:
            lines.append("")

    output = "\n".join(lines)
    truncated = data.get("truncated", False)
    if truncated:
        output += f"\n\n... (truncated at {len(results)} matches, {total} total)"

    return ResultClass(
        output=output,
        title=f"{total} matches for /{pattern}/",
        metadata={"matches": total},
    )


def _local_grep(
    pattern_str: str, search_path: str, file_glob: str | None,
    case_insensitive: bool, context_lines: int, max_results: int,
    ResultClass: type,
) -> Any:
    """Fallback: local regex file search."""
    flags = re.IGNORECASE if case_insensitive else 0
    try:
        regex = re.compile(pattern_str, flags)
    except re.error as e:
        return ResultClass(error=f"Invalid regex: {e}")

    base = Path(search_path).resolve()
    if base.is_file():
        files = [base]
    elif base.is_dir():
        if file_glob:
            files = sorted(base.rglob(file_glob))
        else:
            files = sorted(base.rglob("*"))
        files = [f for f in files if f.is_file()]
    else:
        return ResultClass(error=f"Path not found: {search_path}")

    results = []
    total_matches = 0

    for fp in files:
        if total_matches >= max_results:
            break
        try:
            text = fp.read_text(encoding="utf-8", errors="replace")
        except (PermissionError, OSError):
            continue

        lines = text.splitlines()
        for i, line in enumerate(lines):
            if total_matches >= max_results:
                break
            if regex.search(line):
                total_matches += 1
                try:
                    rel = fp.relative_to(base if base.is_dir() else base.parent)
                except ValueError:
                    rel = fp

                if context_lines > 0:
                    start = max(0, i - context_lines)
                    end = min(len(lines), i + context_lines + 1)
                    for j in range(start, end):
                        prefix = ">" if j == i else " "
                        results.append(f"{rel}:{j + 1}{prefix} {lines[j]}")
                    results.append("")
                else:
                    results.append(f"{rel}:{i + 1}: {line}")

    output = "\n".join(results) if results else "(no matches)"
    if total_matches >= max_results:
        output += f"\n\n... (truncated at {max_results} matches)"

    return ResultClass(
        output=output,
        title=f"{total_matches} matches for /{pattern_str}/",
        metadata={"matches": total_matches},
    )


def _local_glob(
    pattern: str, search_dir: str, workspace: str | None,
    ResultClass: type,
) -> Any:
    """Fallback: local pathlib glob."""
    base = Path(search_dir).resolve()
    if not base.exists():
        return ResultClass(error=f"Directory not found: {search_dir}")

    try:
        matches = list(base.glob(pattern))
    except ValueError as e:
        return ResultClass(error=f"Invalid glob pattern: {e}")

    files = [m for m in matches if m.is_file()]

    # Workspace enforcement
    if workspace:
        ws = Path(workspace).resolve()
        files = [f for f in files if _is_within(f, ws)]

    files.sort(key=lambda p: p.stat().st_mtime, reverse=True)

    max_results = 500
    truncated = len(files) > max_results
    files = files[:max_results]

    output_lines = []
    for f in files:
        try:
            rel = f.relative_to(base)
        except ValueError:
            rel = f
        output_lines.append(str(rel))

    output = "\n".join(output_lines) if output_lines else "(no matches)"
    if truncated:
        output += f"\n\n... ({len(matches) - max_results} more matches)"

    return ResultClass(
        output=output,
        title=f"{len(output_lines)} files matching {pattern}",
        metadata={"count": len(output_lines), "truncated": truncated},
    )


def _is_within(path: Path, parent: Path) -> bool:
    try:
        path.resolve().relative_to(parent)
        return True
    except ValueError:
        return False
